#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
演示内容的附件生成器 —— `scripts/demo/assets/` 里那几个文件就是这个脚本产出的。

为什么要committed 的生成器：附件必须跟着 git 走（`LOCAL_STORAGE_DIR` 不在版本库里，
迁移时也带不走数据库），所以真正入库的是 `scripts/demo/assets/` 下的二进制文件；
但二进制文件一旦想改一个字就没法 review，于是把「怎么画出来的」也留在仓库里。

    python3 scripts/demo/make-assets.py          # 重新生成全部资产
    python3 scripts/demo/make-assets.py --check  # 只报告缺失/体积，不写文件

依赖：Pillow（图片 / PDF）、python-pptx（幻灯片）、openpyxl（表格）。
字体：脚本按 FONT_CANDIDATES 顺序找一款中文字体；找不到就明确报错，
绝不静默生成一堆豆腐块。全部资产加起来控制在几百 KB。

注意：PDF 是「图片版」PDF（每页一张渲染好的位图）。这是刻意的取舍——仓库里没有
reportlab，也没有 LibreOffice，而手写一个内嵌中文字体子集的矢量 PDF 不值得为一份
演示材料付出。浏览器自带的 PDF 阅读器照样能翻页 / 缩放 / 选页。

PDF 由本文件里的 `write_pdf` 自己拼字节，而不是 `Image.save(..., "PDF")`：Pillow 的
PDF 分支把 RGB 页面按 JPEG 存，一份三页的纯文字报告要 300 KB 起；换成「自适应调色板
+ FlateDecode」是无损的，还只要三分之一。仓库里躺着的二进制越小越好。
"""

from __future__ import annotations

import argparse
import math
import os
import sys
import zlib

ROOT = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(ROOT, "assets")

FONT_CANDIDATES = [
    "/System/Library/Fonts/Hiragino Sans GB.ttc",
    "/System/Library/Fonts/STHeiti Medium.ttc",
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
    "C:/Windows/Fonts/msyh.ttc",
]

try:
    from PIL import Image, ImageDraw, ImageFilter, ImageFont
except ImportError:  # pragma: no cover - developer tooling
    print("需要 Pillow：pip3 install pillow", file=sys.stderr)
    raise SystemExit(2)


def font_path() -> str:
    for p in FONT_CANDIDATES:
        if os.path.exists(p):
            try:
                ImageFont.truetype(p, 20)
                return p
            except OSError:
                continue
    print(
        "找不到可用的中文字体。请在 FONT_CANDIDATES 里加上本机的一款中文 TTF/TTC。",
        file=sys.stderr,
    )
    raise SystemExit(2)


FONT = font_path()


def f(size: int):
    return ImageFont.truetype(FONT, size)


def save_png(img, path: str, colors: int = 48) -> None:
    """扁平色的示意图用自适应调色板存 —— 视觉上无差别，体积是真彩 PNG 的三分之一。"""
    img.convert("P", palette=Image.ADAPTIVE, colors=colors).save(path, "PNG", optimize=True)


INK = (24, 24, 27)
MUTED = (113, 113, 122)
LINE = (228, 228, 231)
PAPER = (255, 255, 255)


# ── 封面 ────────────────────────────────────────────────────────────────────


def zone_cover(path: str, base: tuple[int, int, int], accent: tuple[int, int, int],
               title: str, sub: str) -> None:
    """
    版块头图：深色底 + 细网格 + 一条贯穿全宽的色块带。

    页头把它裁成一条很扁的横条（1600×600 → 约 1200×145），所以视觉元素必须
    **横向铺满**：只在右上角放几个块的话，多数裁切窗口里就是一片空黑。
    """
    w, h = 1600, 600
    img = Image.new("RGB", (w, h), base)
    d = ImageDraw.Draw(img, "RGBA")

    for y in range(h):  # 上亮下暗的竖向渐变，逐行画够快也够小
        t = y / max(1, h - 1)
        d.line([(0, y), (w, y)], fill=tuple(int(c * (1 - 0.35 * t)) for c in base))

    step = 48
    for x in range(0, w, step):
        d.line([(x, 0), (x, h)], fill=(255, 255, 255, 14))
    for y in range(0, h, step):
        d.line([(0, y), (w, y)], fill=(255, 255, 255, 14))

    # 一排高低错落的圆角块，横跨整幅：无论裁到哪一段都有东西。
    heights = [150, 230, 110, 280, 170, 320, 130, 210, 260, 140]
    bw, gap = 118, 40
    x = 70
    for i, bh in enumerate(heights):
        y = (h - bh) // 2 + (28 if i % 2 else -28)
        alpha = 26 + (i % 4) * 16
        d.rounded_rectangle([x, y, x + bw, y + bh], radius=14, fill=(*accent, alpha))
        x += bw + gap

    # 底部压一层暗色，保证叠在上面的白字任何时候都读得清。
    d.rectangle([0, h - 150, w, h], fill=(0, 0, 0, 90))
    d.line([(64, h - 96), (184, h - 96)], fill=(*accent, 230), width=3)
    d.text((64, h - 80), title, font=f(30), fill=(244, 244, 245))
    d.text((64, h - 40), sub, font=f(17), fill=(200, 200, 210))
    img.save(path, "JPEG", quality=78, optimize=True)


def post_cover(path: str, accent: tuple[int, int, int], tint: tuple[int, int, int],
               kicker: str, title: str) -> None:
    """
    帖子封面：浅底 + 一块大色块 + 墨色标题。

    列表里它只有指甲盖大，深色抽象图在那个尺寸下就是一团黑；浅底加一块饱和色
    才认得出来是哪一篇。
    """
    w, h = 1200, 630
    img = Image.new("RGB", (w, h), (250, 250, 249))
    d = ImageDraw.Draw(img, "RGBA")

    d.rectangle([0, 0, w, h], fill=(*tint, 255))
    for x in range(0, w, 40):  # 极淡的竖纹，缩略图里只剩一点肌理
        d.line([(x, 0), (x, h)], fill=(255, 255, 255, 60))

    # 右侧的大色块 + 两条错位的窄条 —— 缩略图里就是「一块颜色」，够用了
    d.rounded_rectangle([w - 470, 96, w - 96, 470], radius=28, fill=(*accent, 235))
    d.rounded_rectangle([w - 560, 210, w - 430, 300], radius=18, fill=(*accent, 130))
    d.rounded_rectangle([w - 620, 330, w - 520, 400], radius=14, fill=(*accent, 80))

    d.line([(96, 150), (216, 150)], fill=(*accent, 255), width=6)
    d.text((96, 186), kicker, font=f(28), fill=(90, 90, 100))
    y = 250
    for line in title.split("\n"):
        d.text((96, y), line, font=f(58), fill=(24, 24, 27))
        y += 82
    img.save(path, "JPEG", quality=80, optimize=True)


# ── 正文插图 ────────────────────────────────────────────────────────────────


def fig_pipeline(path: str) -> None:
    """推理服务分层示意图：请求 → 调度 → 执行 → 返回。"""
    w, h = 1400, 720
    img = Image.new("RGB", (w, h), PAPER)
    d = ImageDraw.Draw(img)

    d.text((48, 40), "图 1  推理服务分层示意（示例结构，非生产拓扑）", font=f(26), fill=INK)
    d.line([(48, 88), (w - 48, 88)], fill=LINE, width=2)

    rows = [
        ("接入层", ["HTTP / gRPC 网关", "鉴权与配额", "请求归一化"], (219, 234, 254), (30, 64, 175)),
        ("调度层", ["连续批处理队列", "序列长度分桶", "抢占与重排"], (220, 252, 231), (22, 101, 52)),
        ("执行层", ["算子融合内核", "KV Cache 分页", "INT4 权重加载"], (254, 243, 199), (146, 64, 14)),
        ("观测层", ["逐阶段计时", "显存水位采样", "回归基线比对"], (243, 244, 246), (63, 63, 70)),
    ]
    top = 124
    for i, (name, items, bg, fg) in enumerate(rows):
        y = top + i * 138
        d.rounded_rectangle([48, y, w - 48, y + 108], radius=14, fill=bg)
        d.text((80, y + 34), name, font=f(30), fill=fg)
        for j, it in enumerate(items):
            x = 300 + j * 360
            d.rounded_rectangle([x, y + 24, x + 320, y + 84], radius=10, outline=fg, width=2, fill=PAPER)
            d.text((x + 20, y + 42), it, font=f(21), fill=INK)
        if i < len(rows) - 1:
            cx = w // 2
            d.line([(cx, y + 108), (cx, y + 138)], fill=MUTED, width=3)
            d.polygon([(cx - 8, y + 130), (cx + 8, y + 130), (cx, y + 140)], fill=MUTED)

    d.text((48, h - 44), "每一层都单独计时，任何一层的改动都能在同一张基线上比较。", font=f(20), fill=MUTED)
    save_png(img, path)


def fig_latency(path: str) -> None:
    """两条折线：批大小 vs P50 / P99 延迟（示例数据）。"""
    w, h = 1400, 760
    img = Image.new("RGB", (w, h), PAPER)
    d = ImageDraw.Draw(img)

    d.text((48, 36), "图 2  批大小与端到端延迟（示例数据，仅用于说明形状）", font=f(26), fill=INK)

    left, right, top, bottom = 120, w - 60, 110, h - 110
    batches = [1, 2, 4, 8, 16, 32, 64]
    p50 = [42, 46, 55, 71, 104, 168, 296]
    p99 = [88, 94, 112, 148, 214, 352, 640]
    ymax = 700

    for i in range(6):
        y = bottom - (bottom - top) * i / 5
        d.line([(left, y), (right, y)], fill=LINE, width=1)
        d.text((56, y - 12), str(int(ymax * i / 5)), font=f(18), fill=MUTED)
    d.line([(left, top), (left, bottom)], fill=(161, 161, 170), width=2)
    d.line([(left, bottom), (right, bottom)], fill=(161, 161, 170), width=2)

    def pt(i: int, v: int) -> tuple[float, float]:
        x = left + (right - left) * i / (len(batches) - 1)
        y = bottom - (bottom - top) * min(v, ymax) / ymax
        return x, y

    # 图例固定在左上的空白区，不跟着线走 —— 贴着线端放会压到曲线上。
    legend_y = top + 10
    for series, color, label in ((p50, (37, 99, 235), "P50 延迟 (ms)"), (p99, (220, 38, 38), "P99 延迟 (ms)")):
        pts = [pt(i, v) for i, v in enumerate(series)]
        d.line(pts, fill=color, width=4, joint="curve")
        for (x, y), v in zip(pts, series):
            d.ellipse([x - 6, y - 6, x + 6, y + 6], fill=color)
            d.text((x - 16, y - 34), f"{v}", font=f(17), fill=color)
        d.line([(left + 40, legend_y + 12), (left + 90, legend_y + 12)], fill=color, width=4)
        d.text((left + 102, legend_y), label, font=f(21), fill=color)
        legend_y += 36

    for i, b in enumerate(batches):
        x, _ = pt(i, 0)
        d.text((x - 10, bottom + 14), str(b), font=f(19), fill=MUTED)
    d.text((right - 74, bottom + 48), "批大小", font=f(21), fill=INK)
    d.text((48, h - 40), "吞吐随批大小增长，尾延迟增长更快 —— 预算怎么定取决于哪一条先撞线。",
           font=f(20), fill=MUTED)
    save_png(img, path)


def fig_frame_budget(path: str) -> None:
    """帧预算堆叠条：一帧 16.6ms 被各阶段吃掉多少（示例数据）。"""
    w, h = 1080, 600
    img = Image.new("RGB", (w, h), PAPER)
    d = ImageDraw.Draw(img)
    d.text((48, 36), "图 1  一帧 16.6 ms 的去向（示例数据）", font=f(26), fill=INK)

    stages = [
        ("几何提交", 2.4, (191, 219, 254)),
        ("阴影通道", 4.1, (167, 243, 208)),
        ("光照合成", 5.2, (253, 230, 138)),
        ("后处理", 2.6, (233, 213, 255)),
        ("呈现等待", 2.3, (228, 228, 231)),
    ]
    total = sum(s[1] for s in stages)
    left, right = 80, w - 80
    y0, y1 = 150, 250
    x = left
    for name, val, color in stages:
        seg = (right - left) * val / total
        d.rectangle([x, y0, x + seg, y1], fill=color, outline=PAPER, width=3)
        d.text((x + 12, y0 + 34), name, font=f(21), fill=INK)
        x += seg

    d.text((80, 282), f"合计 {total:.1f} ms / 预算 16.6 ms", font=f(24), fill=INK)

    ly = 360
    for name, val, color in stages:
        d.rectangle([80, ly, 116, ly + 26], fill=color)
        d.text((132, ly + 1), f"{name}    {val:.1f} ms    {val / total * 100:.0f}%", font=f(22), fill=INK)
        ly += 44
    save_png(img, path)


def fig_queue(path: str) -> None:
    """堆叠柱：不同并发下一次请求的时间构成（示例数据）。"""
    w, h = 1160, 660
    img = Image.new("RGB", (w, h), PAPER)
    d = ImageDraw.Draw(img)
    d.text((48, 36), "图 1  一次请求的时间构成随并发的变化（示例数据）", font=f(26), fill=INK)

    groups = [
        ("并发 4", [(28, "排队等待"), (34, "预填充"), (52, "解码")]),
        ("并发 16", [(61, "排队等待"), (36, "预填充"), (58, "解码")]),
        ("并发 32", [(118, "排队等待"), (38, "预填充"), (61, "解码")]),
        ("并发 64", [(236, "排队等待"), (41, "预填充"), (64, "解码")]),
    ]
    colors = [(148, 163, 184), (96, 165, 250), (52, 211, 153)]
    left, base_y = 150, 470
    scale = 1.05
    bar_w, gap = 120, 100
    for gi, (label, parts) in enumerate(groups):
        x = left + gi * (bar_w + gap)
        y = base_y
        for pi, (val, _name) in enumerate(parts):
            hgt = val * scale
            d.rectangle([x, y - hgt, x + bar_w, y], fill=colors[pi], outline=PAPER, width=2)
            if hgt > 26:
                d.text((x + 10, y - hgt / 2 - 12), f"{val}", font=f(19), fill=(24, 24, 27))
            y -= hgt
        total = sum(v for v, _ in parts)
        d.text((x + 18, base_y + 14), label, font=f(21), fill=INK)
        d.text((x + 10, y - 34), f"{total} ms", font=f(21), fill=INK)
    d.line([(left - 30, base_y), (w - 60, base_y)], fill=(161, 161, 170), width=2)

    lx = 150
    for (_val, name), color in zip(groups[0][1], colors):
        d.rectangle([lx, 540, lx + 32, 564], fill=color)
        d.text((lx + 44, 540), name, font=f(21), fill=INK)
        lx += 240
    d.text((48, h - 44), "并发一高，多出来的时间几乎全落在排队上 —— 计算部分基本没变。",
           font=f(20), fill=MUTED)
    save_png(img, path)


def checklist_md(path: str) -> None:
    """一份纯文本附件：预览面板会把 md/txt/csv/json 直接排进 <pre>。"""
    text = """# 帧预算排查清单（演示用）

> 这份清单只是演示附件，用来说明「文本类附件可以直接在阅读面板里读」。

## 开始之前

- [ ] 锁定机型与系统版本，记录在实验条目里
- [ ] 关闭动态调频，避免两次跑分不可比
- [ ] 确认采样窗口 >= 600 帧，丢掉前 60 帧

## 逐阶段

| 阶段 | 关注指标 | 常见坑 |
| --- | --- | --- |
| 几何提交 | Draw Call 数 | 合批被材质切换打断 |
| 阴影通道 | 级联数 / 分辨率 | 远处级联分辨率给太高 |
| 光照合成 | 带宽占用 | 中间纹理格式没降精度 |
| 后处理 | 通道数 | 相邻通道没合并 |
| 呈现等待 | 掉帧分布 | 只看均值会漏掉抖动 |

## 收尾

- [ ] 把这一轮的数字写回版块 Wiki 的《实验规范》
- [ ] 有回归的项单独开一条实验记录，别塞进同一篇
"""
    with open(path, "w", encoding="utf-8") as fh:
        fh.write(text)


# ── PDF（图片版） ───────────────────────────────────────────────────────────

PDF_PALETTE_COLORS = 32


def write_pdf(path: str, pages: list[Image.Image], dpi: float) -> None:
    """
    把若干张位图写成一个多页 PDF：每页一个 Indexed/FlateDecode 的 XObject。

    结构是最小可用的那一套（Catalog → Pages → 每页 Page + Contents + Image），
    交叉引用表按实际偏移量写死。之所以不用 Pillow 的 PDF 分支，见文件头的说明。
    """
    objs: list[bytes] = []

    def add(body: bytes) -> int:
        objs.append(body)
        return len(objs)

    add(b"")  # 1: Catalog（占位，最后回填）
    add(b"")  # 2: Pages
    kids: list[int] = []

    for im in pages:
        p = im.convert("P", palette=Image.ADAPTIVE, colors=PDF_PALETTE_COLORS)
        palette = bytes(p.getpalette()[: PDF_PALETTE_COLORS * 3])
        data = zlib.compress(p.tobytes(), 9)
        pw, ph = p.size
        w, h = pw / dpi * 72.0, ph / dpi * 72.0
        img = add(
            b"<</Type/XObject/Subtype/Image/Width %d/Height %d"
            b"/ColorSpace[/Indexed/DeviceRGB %d<%s>]/BitsPerComponent 8"
            b"/Filter/FlateDecode/Length %d>>stream\n%s\nendstream"
            % (pw, ph, PDF_PALETTE_COLORS - 1, palette.hex().encode(), len(data), data)
        )
        stream = zlib.compress(b"q %.2f 0 0 %.2f 0 0 cm /Im0 Do Q" % (w, h), 9)
        content = add(b"<</Filter/FlateDecode/Length %d>>stream\n%s\nendstream" % (len(stream), stream))
        kids.append(add(
            b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 %.2f %.2f]"
            b"/Resources<</XObject<</Im0 %d 0 R>>>>/Contents %d 0 R>>" % (w, h, img, content)
        ))

    objs[0] = b"<</Type/Catalog/Pages 2 0 R>>"
    objs[1] = b"<</Type/Pages/Kids[%s]/Count %d>>" % (b" ".join(b"%d 0 R" % k for k in kids), len(kids))

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for i, body in enumerate(objs, start=1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % i + body + b"\nendobj\n"
    xref = len(out)
    out += b"xref\n0 %d\n0000000000 65535 f \n" % (len(objs) + 1)
    for off in offsets:
        out += b"%010d 00000 n \n" % off
    out += b"trailer<</Size %d/Root 1 0 R>>\nstartxref\n%d\n%%%%EOF\n" % (len(objs) + 1, xref)
    with open(path, "wb") as fh:
        fh.write(bytes(out))


def pdf_report(path: str) -> None:
    """三页技术报告：封面 / 方法与配置表 / 结论。"""
    W, H = 1240, 1754  # A4 @150dpi

    def page() -> tuple[Image.Image, ImageDraw.ImageDraw]:
        im = Image.new("RGB", (W, H), PAPER)
        return im, ImageDraw.Draw(im)

    # ── 封面
    p1, d = page()
    d.rectangle([0, 0, W, 220], fill=(24, 24, 27))
    d.text((80, 84), "异构推理加速  ·  阶段评测报告", font=f(44), fill=(250, 250, 250))
    d.text((80, 148), "温哥华研究所 / Computing Data Application Acceleration Laboratory",
           font=f(20), fill=(180, 180, 190))
    d.text((80, 320), "算子融合与 INT4 量化的联合评测", font=f(52), fill=INK)
    d.text((80, 410), "演示材料 · 数据为构造样例，不代表任何真实产品结论", font=f(24), fill=MUTED)
    d.line([(80, 470), (W - 80, 470)], fill=LINE, width=2)
    for i, line in enumerate([
        "1  评测目标与范围",
        "2  基线与硬件配置",
        "3  单点改动的收益拆解",
        "4  组合后的相互作用",
        "5  结论与下一步",
    ]):
        d.text((80, 540 + i * 56), line, font=f(28), fill=INK)
    d.text((80, H - 120), "本文件由 scripts/demo/make-assets.py 生成，仅用于社区功能演示。",
           font=f(20), fill=MUTED)

    # ── 配置表
    p2, d = page()
    d.text((80, 90), "2  基线与硬件配置", font=f(40), fill=INK)
    d.line([(80, 160), (W - 80, 160)], fill=LINE, width=2)
    rows = [
        ("项目", "取值", "说明"),
        ("模型规模", "7B / 稠密", "权重 FP16 基线"),
        ("上下文", "4096", "输入 1024 / 输出 128"),
        ("批处理", "连续批处理", "最大并发 32"),
        ("量化", "INT4 (权重)", "激活保持 FP16"),
        ("算子融合", "Attention + MLP", "融合后减少一次显存往返"),
        ("重复次数", "5 轮取中位数", "丢弃首轮预热"),
    ]
    y = 200
    for i, (a, b, c) in enumerate(rows):
        bg = (244, 244, 245) if i == 0 else PAPER
        d.rectangle([80, y, W - 80, y + 66], fill=bg, outline=LINE, width=1)
        fnt = f(24)
        d.text((100, y + 20), a, font=fnt, fill=INK)
        d.text((420, y + 20), b, font=fnt, fill=INK)
        d.text((700, y + 20), c, font=f(22), fill=MUTED if i else INK)
        y += 66
    d.text((80, y + 50), "所有数字都是为了演示排版而构造的，请勿引用。", font=f(22), fill=MUTED)

    # ── 结论
    p3, d = page()
    d.text((80, 90), "5  结论与下一步", font=f(40), fill=INK)
    d.line([(80, 160), (W - 80, 160)], fill=LINE, width=2)
    bullets = [
        "两项改动单独都有收益，叠加后并不是简单相加：",
        "显存带宽先于算力成为瓶颈，融合的收益随批大小衰减。",
        "尾延迟（P99）对批大小比吞吐敏感得多，",
        "把批大小当成唯一旋钮会先撞上延迟预算。",
        "下一步：把逐阶段计时接入常态化回归，",
        "让每次内核改动都能落在同一张基线上比较。",
    ]
    for i, b in enumerate(bullets):
        d.text((80, 230 + i * 58), b, font=f(28), fill=INK)
    d.rounded_rectangle([80, 640, W - 80, 860], radius=16, outline=LINE, width=2)
    d.text((110, 680), "遗留问题", font=f(30), fill=INK)
    d.text((110, 740), "· 低批场景下融合内核的收益尚未稳定复现", font=f(24), fill=INK)
    d.text((110, 790), "· 量化误差在长上下文下的累积仍缺少长跑数据", font=f(24), fill=INK)

    write_pdf(path, [p1, p2, p3], dpi=150.0)


# ── PPTX / XLSX ─────────────────────────────────────────────────────────────


def pptx_deck(path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def slide(title: str, bullets: list[str]) -> None:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        s.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
        body = s.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0x27, 0x27, 0x2A)

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "批处理调度实验记录"
    cover.placeholders[1].text = "连续批处理 vs 静态批处理 · 演示材料（构造数据）"

    slide("我们在比什么", [
        "同一模型、同一硬件、同一输入分布",
        "只切换调度策略：静态批处理 / 连续批处理",
        "关注三个数：吞吐、P50、P99",
    ])
    slide("怎么测", [
        "固定 5 轮取中位数，丢弃首轮预热",
        "输入 1024 token / 输出 128 token",
        "每一阶段单独计时，避免把等待算进计算",
    ])
    slide("观察到的形状", [
        "吞吐随并发上升，但 P99 上升得更快",
        "静态批处理在低并发下反而更稳",
        "队列等待在高并发下成为主要成分",
    ])
    slide("下一步", [
        "把逐阶段计时接入常态化回归",
        "补一次长跑，观察显存水位是否漂移",
        "把结论写进版块 Wiki 的实验规范",
    ])
    prs.save(path)


BENCH_ROWS = [
    ["FP16 基线", 1, 512, 42, 88, 15.8, "对照组"],
    ["FP16 基线", 8, 2180, 71, 148, 16.4, ""],
    ["FP16 基线", 32, 5120, 168, 352, 18.9, ""],
    ["+ 算子融合", 1, 583, 38, 80, 15.6, "减少一次显存往返"],
    ["+ 算子融合", 8, 2560, 63, 133, 16.1, ""],
    ["+ 算子融合", 32, 5780, 152, 318, 18.4, ""],
    ["+ INT4 量化", 1, 641, 34, 74, 9.2, "权重 INT4 / 激活 FP16"],
    ["+ INT4 量化", 8, 2910, 57, 126, 9.9, ""],
    ["+ INT4 量化", 32, 6340, 141, 305, 12.1, ""],
    ["融合 + INT4", 1, 688, 32, 70, 9.1, "非简单相加"],
    ["融合 + INT4", 8, 3120, 54, 119, 9.7, ""],
    ["融合 + INT4", 32, 6620, 136, 297, 11.8, "带宽先撞线"],
]
BENCH_HEADER = ["配置", "批大小", "吞吐 (tok/s)", "P50 (ms)", "P99 (ms)", "显存峰值 (GB)", "备注"]


def csv_bench(path: str) -> None:
    """
    同一份数据的 CSV。

    存在的理由很实际：xlsx 的面板预览要靠 LibreOffice 转 PDF，没装就只能显示
    「无法预览」；csv/md/txt/json 是面板自己排的，任何机器上都能直接读。演示里
    两份都挂上，既展示了办公文档转换，也保证至少有一份一定打得开。
    """
    import csv as _csv

    with open(path, "w", encoding="utf-8-sig", newline="") as fh:
        w = _csv.writer(fh)
        w.writerow(BENCH_HEADER)
        w.writerows(BENCH_ROWS)
        w.writerow([])
        w.writerow(["# 数据为构造样例，仅用于演示，请勿引用。"])


def xlsx_bench(path: str) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = "评测矩阵"

    ws.append(BENCH_HEADER)
    for r in BENCH_ROWS:
        ws.append(r)

    fill = PatternFill("solid", fgColor="18181B")
    for c in ws[1]:
        c.font = Font(bold=True, color="FAFAFA", size=11)
        c.fill = fill
        c.alignment = Alignment(horizontal="center")
    widths = [16, 10, 15, 11, 11, 16, 26]
    for i, wdt in enumerate(widths, start=1):
        ws.column_dimensions[chr(64 + i)].width = wdt
    ws.freeze_panes = "A2"

    note = wb.create_sheet("说明")
    note["A1"] = "数据为构造样例，仅用于演示社区的附件预览与表格渲染，请勿引用。"
    note.column_dimensions["A"].width = 70
    wb.save(path)


# ── 活动封面 / 讲师头像 ──────────────────────────────────────────────────────
#
# 跟版块头图、帖子封面不同，活动封面**一个字都不印**。理由是渲染方式：列表卡把
# 它裁成 96×96 的方图，详情页按 2:1 铺开 —— 任何烤进图里的字在方图裁切下都会被
# 切掉半行。`public/labs/` 的研究所封面早就是这个结论（见 scripts/make-lab-covers.py
# 的抬头），这里照同一套做：抽象几何、无文字、每张一个低饱和主色。
#
# 主色**按城市**取，而且直接复用研究所封面那六个色 —— 于是温哥华的活动跟温哥华
# 研究所那块砖是同一个色系，站点读起来是一套东西，而不是各画各的。

EV_GROUND = (24, 26, 31)

# 城市 → 主色（前六个与 scripts/make-lab-covers.py 的 CITIES 一一对应）。
EVENT_HUES: dict[str, tuple[int, int, int]] = {
    "vancouver": (72, 138, 156),   # 冷青 — 海与山
    "toronto": (78, 112, 184),     # 蓝 — 天际线
    "ottawa": (158, 94, 94),       # 暗红 — 议会大厦的石头
    "waterloo": (104, 142, 96),    # 绿 — 电路格
    "edmonton": (128, 104, 174),   # 紫 — 极光
    "montreal": (172, 132, 82),    # 琥珀 — 桥
    "china": (166, 106, 106),      # 砖红 — 第七个城市选项
}


def _wash(img, hue: tuple[int, int, int], w: int, h: int, strength: float = 0.9) -> None:
    """墨底之上一层自上而下的主色渐变（与研究所封面同一手法）。"""
    d = ImageDraw.Draw(img)
    for y in range(h):
        t = (1 - y / h) ** 1.3 * strength
        d.line([(0, y), (w, y)], fill=tuple(
            int(EV_GROUND[i] + (hue[i] - EV_GROUND[i]) * t) for i in range(3)
        ))


def _ev_waves(d, hue, w, h):
    """层叠的波带 —— 给会议/峰会。"""
    for i in range(7):
        pts = []
        base = h * (0.30 + i * 0.085)
        for x in range(0, w + 1, 16):
            pts.append((x, base + math.sin(x / 150 + i * 0.8) * (26 + i * 4)))
        d.line(pts, fill=tuple(int(c * (0.9 - i * 0.10)) for c in hue), width=max(1, 8 - i))


def _ev_orbits(d, hue, w, h):
    """同心弧 + 节点 —— 给分享/讲座。"""
    cx, cy = w * 0.72, h * 0.52
    for i in range(6):
        r = 70 + i * 62
        d.arc([cx - r, cy - r, cx + r, cy + r], 150, 390,
              fill=tuple(int(c * (0.85 - i * 0.09)) for c in hue), width=3)
        a = math.radians(200 + i * 33)
        d.ellipse([cx + math.cos(a) * r - 7, cy + math.sin(a) * r - 7,
                   cx + math.cos(a) * r + 7, cy + math.sin(a) * r + 7], fill=hue)


def _ev_lattice(d, hue, w, h):
    """点阵格 —— 给研讨会/工作坊。"""
    step = 62
    for gy in range(1, h // step + 1):
        for gx in range(1, w // step + 1):
            x, y = gx * step, gy * step
            if (gx * 7 + gy * 5) % 4:
                d.line([(x, y), (x + step, y)], fill=tuple(int(c * 0.34) for c in hue), width=1)
            if (gx * 3 + gy * 11) % 5:
                d.line([(x, y), (x, y + step)], fill=tuple(int(c * 0.34) for c in hue), width=1)
            if (gx + gy) % 3 == 0:
                d.ellipse([x - 4, y - 4, x + 4, y + 4], fill=hue)


def _ev_stack(d, hue, w, h):
    """错落的圆角条 —— 像一张日程表。"""
    y = h * 0.20
    for i in range(7):
        x0 = 90 + (i % 3) * 130
        bw = 240 + (i % 4) * 150
        d.rounded_rectangle([x0, y, x0 + bw, y + 34], radius=17,
                            fill=tuple(int(c * (0.95 - i * 0.09)) for c in hue))
        y += 52


def _ev_radial(d, hue, w, h):
    """放射线 —— 给线上活动。"""
    cx, cy = w * 0.5, h * 1.05
    for i in range(26):
        a = math.radians(190 + i * 6.2)
        r = h * (1.15 if i % 3 else 1.35)
        d.line([(cx, cy), (cx + math.cos(a) * r, cy + math.sin(a) * r)],
               fill=tuple(int(c * (0.28 + 0.5 * (i % 4) / 3)) for c in hue), width=2)
    for r in (h * 0.45, h * 0.72, h * 0.99):
        d.arc([cx - r, cy - r, cx + r, cy + r], 180, 360,
              fill=tuple(int(c * 0.8) for c in hue), width=3)


def _ev_contours(d, hue, w, h):
    """嵌套等高线 —— 给多日/大型活动。"""
    cx, cy = w * 0.38, h * 0.5
    for i in range(9):
        pts = []
        for k in range(0, 361, 6):
            a = math.radians(k)
            rr = (60 + i * 42) * (1 + 0.16 * math.sin(a * 3 + i * 0.7))
            pts.append((cx + math.cos(a) * rr * 1.35, cy + math.sin(a) * rr * 0.72))
        d.line(pts + [pts[0]], fill=tuple(int(c * (0.9 - i * 0.075)) for c in hue), width=2)


def _ev_mesh(d, hue, w, h):
    """三角网格 —— 给图形/渲染方向。"""
    step = 96
    rows = h // step + 2
    cols = w // step + 2
    for r in range(rows):
        for c in range(cols):
            x = c * step + (step // 2 if r % 2 else 0) - step
            y = r * step * 0.86 - step
            for dx, dy in ((step, 0), (step // 2, step * 0.86), (-step // 2, step * 0.86)):
                if (r * 5 + c * 3 + int(dx)) % 4:
                    d.line([(x, y), (x + dx, y + dy)],
                           fill=tuple(int(c2 * (0.26 + 0.42 * ((r + c) % 3) / 2)) for c2 in hue), width=1)
            if (r * 3 + c * 7) % 5 == 0:
                d.ellipse([x - 5, y - 5, x + 5, y + 5], fill=hue)


EVENT_MOTIFS = {
    "waves": _ev_waves,
    "orbits": _ev_orbits,
    "lattice": _ev_lattice,
    "stack": _ev_stack,
    "radial": _ev_radial,
    "contours": _ev_contours,
    "mesh": _ev_mesh,
}


def event_cover(path: str, city: str, motif: str) -> None:
    """2:1 的活动封面（详情页按 2:1 铺，列表卡再裁成方图）。"""
    w, h = 1200, 600
    img = Image.new("RGB", (w, h), EV_GROUND)
    _wash(img, EVENT_HUES[city], w, h)
    # 直接画在洗好的底上：先画到透明层再合成，出来的就是一块近黑的方块
    # （make-lab-covers.py 第一版踩过这个坑）。
    EVENT_MOTIFS[motif](ImageDraw.Draw(img), tuple(min(255, int(c * 1.45)) for c in EVENT_HUES[city]), w, h)
    img = img.filter(ImageFilter.SMOOTH)
    # 四角压暗一点，卡片缩略图里边缘不至于糊在背景上。
    v = Image.new("L", (w, h), 255)
    ImageDraw.Draw(v).rectangle([0, 0, w, h], outline=150, width=90)
    img = Image.composite(img, Image.new("RGB", (w, h), EV_GROUND), v.filter(ImageFilter.GaussianBlur(60)))
    img.save(path, "JPEG", quality=78, optimize=True)


def speaker_avatar(path: str, hue: tuple[int, int, int], seed: int) -> None:
    """
    讲师头像占位图：抽象剪影，**不是**假照片，也没有文字。

    详情页把它渲染成 96–112 px 的方块。刻意只给一部分讲师配图，剩下的走
    Avatar 的首字母兜底 —— 两条分支在演示里都要看得见。
    """
    s = 256
    img = Image.new("RGB", (s, s), EV_GROUND)
    _wash(img, hue, s, s, strength=0.75)
    d = ImageDraw.Draw(img, "RGBA")
    light = tuple(min(255, int(c * 1.6)) for c in hue)
    # 头 + 肩：够像一个人，又明显是占位图。
    d.ellipse([s * 0.34, s * 0.20, s * 0.66, s * 0.52], fill=(*light, 235))
    d.pieslice([s * 0.16, s * 0.56, s * 0.84, s * 1.24], 180, 360, fill=(*light, 200))
    # 一条随机相位的细弧，让六张头像各不相同。
    for i in range(3):
        r = s * (0.30 + i * 0.09)
        a0 = (seed * 37 + i * 61) % 360
        d.arc([s / 2 - r, s / 2 - r, s / 2 + r, s / 2 + r], a0, a0 + 70,
              fill=(*light, 120 - i * 30), width=3)
    img.filter(ImageFilter.SMOOTH).save(path, "JPEG", quality=76, optimize=True)



# ── main ────────────────────────────────────────────────────────────────────

# 活动封面：文件名 → (城市主色, 抽象母题)。城市决定颜色，母题决定图形，
# 两个都不印字（列表卡会把它裁成方图）。
EVENT_COVERS: dict[str, tuple[str, str]] = {
    "conf-montreal": ("montreal", "contours"),
    "workshop-vancouver": ("vancouver", "mesh"),
    "hack-vancouver": ("vancouver", "lattice"),
    "seminar-ottawa": ("ottawa", "stack"),
    "openday-waterloo": ("waterloo", "radial"),
    "summit-toronto": ("toronto", "waves"),
    "forum-montreal": ("montreal", "orbits"),
}

# 讲师头像占位图：文件名 → 取色的城市。
SPEAKER_AVATARS: dict[str, str] = {
    "a": "toronto",
    "b": "vancouver",
    "c": "edmonton",
    "d": "waterloo",
    "e": "china",
    "f": "montreal",
}


TARGETS = [
    "zone-cover-inference.jpg",
    "zone-cover-graphics.jpg",
    "post-cover-quant.jpg",
    "post-cover-batching.jpg",
    "post-cover-frame.jpg",
    "fig-pipeline.png",
    "fig-latency.png",
    "fig-frame-budget.png",
    "fig-queue.png",
    "checklist-frame-budget.md",
    "report-inference-eval.pdf",
    "deck-batching.pptx",
    "bench-matrix.xlsx",
    "bench-matrix.csv",
] + [f"event-cover-{n}.jpg" for n in EVENT_COVERS] + [f"speaker-{n}.jpg" for n in SPEAKER_AVATARS]


def build() -> None:
    os.makedirs(OUT, exist_ok=True)
    p = lambda n: os.path.join(OUT, n)  # noqa: E731

    zone_cover(p("zone-cover-inference.jpg"), (23, 32, 51), (96, 165, 250),
               "异构推理加速", "温哥华研究所 · Computing Data Application Acceleration Laboratory")
    zone_cover(p("zone-cover-graphics.jpg"), (28, 25, 44), (167, 139, 250),
               "图形渲染技术", "温哥华研究所 · Graphics Technology Laboratory")
    post_cover(p("post-cover-quant.jpg"), (13, 148, 136), (236, 253, 245),
               "技术报告 · 阶段评测", "算子融合\n× INT4 量化")
    post_cover(p("post-cover-batching.jpg"), (180, 83, 9), (255, 251, 235),
               "实验记录 · 调度", "连续批处理\nvs 静态批处理")
    post_cover(p("post-cover-frame.jpg"), (109, 40, 217), (245, 243, 255),
               "性能剖析 · 渲染管线", "一帧 16.6 ms\n的去向")

    fig_pipeline(p("fig-pipeline.png"))
    fig_latency(p("fig-latency.png"))
    fig_frame_budget(p("fig-frame-budget.png"))
    fig_queue(p("fig-queue.png"))
    checklist_md(p("checklist-frame-budget.md"))
    pdf_report(p("report-inference-eval.pdf"))
    pptx_deck(p("deck-batching.pptx"))
    xlsx_bench(p("bench-matrix.xlsx"))
    csv_bench(p("bench-matrix.csv"))

    for name, (city, motif) in EVENT_COVERS.items():
        event_cover(p(f"event-cover-{name}.jpg"), city, motif)
    for i, (name, city) in enumerate(SPEAKER_AVATARS.items()):
        speaker_avatar(p(f"speaker-{name}.jpg"), EVENT_HUES[city], i + 1)


def report() -> int:
    total = 0
    missing = 0
    for name in TARGETS:
        fp = os.path.join(OUT, name)
        if not os.path.exists(fp):
            print(f"  MISSING  {name}")
            missing += 1
            continue
        size = os.path.getsize(fp)
        total += size
        print(f"  {size / 1024:8.1f} KB  {name}")
    print(f"  {'-' * 30}\n  {total / 1024:8.1f} KB  合计（{len(TARGETS) - missing}/{len(TARGETS)} 个文件）")
    return missing


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="只报告，不写文件")
    args = ap.parse_args()
    if not args.check:
        build()
        print(f"字体：{FONT}")
    raise SystemExit(1 if report() and args.check else 0)
